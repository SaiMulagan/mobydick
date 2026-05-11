"""
Rebuild MobyDicks_Tech_v2.pptx for the technical milestone.

Starts from the existing 16-slide deck (which has both the product pitch
and the technical section), then:
  1. Drops slides 1-6 (the product pitch — already presented).
  2. Reframes the old "PART TWO" slide as a clean technical-opener.
  3. Updates text on the remaining slides to match what's actually shipped
     in the repo (Phase 2 personalization, Reddit ingestion + signal,
     gemini-2.5-flash, SQLite + SQLModel, auto-exclude feedback loop).
  4. Preserves the wine/cream/Georgia-italic styling and the manually
     positioned system-design diagram on what was slide 10.

Run from repo root with the conda env's python.
"""

from copy import deepcopy
from pptx import Presentation

SRC = "/Users/aatishlobo/Desktop/mobydick/MobyDicks_Tech_v2.pptx"
OUT = "/Users/aatishlobo/Desktop/mobydick/MobyDicks_Tech_v2.pptx"


def delete_slide(prs, slide_idx):
    """python-pptx has no public delete; manipulate the XML slide list."""
    rId = prs.slides._sldIdLst[slide_idx].rId  # noqa: N806
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_idx]


def replace_in_runs(slide, mapping):
    """For every run on the slide whose text is exactly a key in `mapping`,
    replace it with the value. Preserves font/size/color/bold/italic."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text in mapping:
                    run.text = mapping[run.text]


def replace_first_run_of_para(slide, old_para_text, new_text):
    """For paragraphs whose concatenated run text equals `old_para_text`,
    put `new_text` in the first run and clear the rest. Used when a
    paragraph is split across multiple runs."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            joined = "".join(r.text for r in para.runs)
            if joined == old_para_text and para.runs:
                para.runs[0].text = new_text
                for r in para.runs[1:]:
                    r.text = ""


def main():
    prs = Presentation(SRC)

    # 1) Drop the product-pitch section (originally slides 1-6).
    #    After each delete, indices shift; always delete index 0.
    for _ in range(6):
        delete_slide(prs, 0)

    # The deck now starts at what was slide 7 ("PART TWO").
    # New indexing: 0..9 (10 slides).

    # ----- Slide 0  (was 7: "PART TWO" opener) -------------------------------
    s = prs.slides[0]
    replace_in_runs(s, {
        "PART  TWO": "TECHNICAL  PRESENTATION",
        "Final Technical Slides": "Moby Dicks · Technical",
        "System design  ·  stack  ·  request flow  ·  demo  ·  repository":
            "Data · models · architecture · testing · demo",
        "The slides that follow contain the final system design and technical artifacts for Milestone 5.":
            "Final system design and technical artifacts for the Advanced MLOps technical milestone.",
    })

    # ----- Slide 1  (was 8: Introduction) ------------------------------------
    s = prs.slides[1]
    replace_in_runs(s, {
        "Extract → fetch → curate. Grounded in a 2.36M-book BigQuery catalog.":
            "Extract → fetch → curate. Grounded in a 2.3M-book BigQuery catalog + Reddit signal.",
        "Ordered by progressive difficulty, each pick justified, ~10 s end-to-end.":
            "Five picks with reasoning, plus per-pick check-ins and a dashboard. ~5 s on a cache hit.",
        "A user submits a topic, a time budget, and a difficulty level. The app returns a five-book curriculum with a written justification per book and an overall reading arc.":
            "A user submits a topic, a time budget, and a difficulty level. The app returns a five-book curriculum with written reasoning and a per-book check-in loop that informs future recommendations.",
    })

    # ----- Slide 2  (was 9: Stack) -------------------------------------------
    s = prs.slides[2]
    replace_in_runs(s, {
        "GCP services + the four required pieces.":
            "GCP services + the four required pieces (Docker · CI/CD · Cloud · MLflow).",
        "2.36M+ Goodreads books; SQL-driven candidate pool.":
            "2.3M Goodreads books + 498 Reddit posts; parameterized SQL pool.",
        "Two structured-JSON calls per request.":
            "Two structured-JSON calls per request; model gemini-2.5-flash.",
        "Single image: FastAPI on :8000, MLflow on :5000.":
            "Container image targeting Cloud Run; build script in repo root.",
        "Nightly cron + manual dispatch run pytest in Ubuntu 24.04.":
            "Nightly cron + workflow_dispatch run pytest on Ubuntu 24.04.",
        "Tracking server inside the container; runs logged per build.":
            "Run tracking stubbed in CI (MagicMock); planned for Phase 3 reranker.",
        "Scale-to-zero service, IAM-authenticated, secrets injected at boot.":
            "Scale-to-zero target; ADC for Vertex AI + BigQuery + Cloud Logging.",
    })

    # ----- Slide 3  (was 10: System Design diagram) --------------------------
    # The diagram is positioned manually; leave it alone. Only minor text
    # cleanups would be possible — skip to avoid breaking the layout.

    # ----- Slide 4  (was 11: Data + Database) --------------------------------
    s = prs.slides[4]
    replace_in_runs(s, {
        "Goodreads catalog in BigQuery, queried at request time.":
            "Three data layers: Goodreads catalog · Reddit signal · per-user state.",
        "books in the table": "books · 498 Reddit posts · 628 mentions",
        "UCSD Goodreads dataset (Wan & McAuley 2018) loaded via load_to_bigquery.sh":
            "UCSD Goodreads dump · r/books + r/suggestmeabook + r/literature via PRAW-free JSON · SQLite for user data",
        "book_id, title, description, num_pages, publication_year":
            "BIGQUERY  ·  Goodreads canonical book + author tables (typed)",
        "average_rating, ratings_count, language_code":
            "BIGQUERY  ·  reddit_posts + reddit_book_mentions (Gemini-extracted)",
        "popular_shelves[ ]   ← shelf-name + count tuples":
            "BIGQUERY  ·  v_reddit_signal view — recommended / asked-similar / weighted",
        "authors[ ]   ← author_id graph for de-dup":
            "SQLITE  ·  user / curriculum / progress / read_book (Phase 2)",
        "1. Filter by language, page range, ratings_count >= 100":
            "1. Filter by language, page range, ratings_count >= 100",
        "2. Strip noise shelves (to-read, favorites, kindle, …)":
            "2. Strip noise shelves; exclude books user has already finished",
        "3. Score by overlap between user keywords and shelf tokens":
            "3. Score by keyword↔shelf-token overlap; tiebreak on Reddit signal",
        "4. Deduplicate by work_id, return top 200, take top 25 to model":
            "4. Dedup by work_id; top 200 → top 25 candidates handed to Gemini",
    })

    # ----- Slide 5  (was 12: Models) -----------------------------------------
    s = prs.slides[5]
    replace_in_runs(s, {
        "Two schema-enforced Gemini Flash calls.":
            "Two schema-enforced Gemini Flash calls + one offline extraction.",
        "Model: gemini-2.0-flash via Vertex AI":
            "Model: gemini-2.5-flash via Vertex AI",
        "Temperature: 0.2  ·  Output: JSON":
            "Temperature: 0.2  ·  response_schema: SurveyFilters",
        "Schema: SurveyFilters (Pydantic)":
            "is_recognized_topic flag short-circuits nonsense input early.",
        'Maps free text ("evening read", "18th-century Russian lit") to keywords, page range, era, and difficulty.':
            'Maps "evening read" / "18th-century Russian lit" → keywords + page-range + era + difficulty.',
        "Why separate: lets us run a deterministic SQL query before the model sees any books — cheaper and grounded.":
            "Why two calls: deterministic SQL runs before the model sees any books — cheaper, grounded, no hallucinations.",
        "Temperature: 0.3  ·  Output: JSON":
            "Temperature: 0.3  ·  response_schema: Curriculum (5 picks)",
        "Schema: Curriculum (week, title, author, reason, overall_arc)":
            "Tight output budget: reason ≤ 180 chars, arc ≤ 250 chars (10 s → 5 s).",
        "Receives top-25 candidates and orders them by progressive difficulty.":
            "Receives the top-25 candidates and orders by progressive difficulty.",
        "No hallucinations: the model can only choose from real book_ids in the pool. Phase 2 will add an MLflow-tracked reranker over check-in data.":
            "Bonus: a third offline Gemini call extracts book mentions from Reddit posts. Phase 3 will train an MLflow-tracked reranker on the check-in data Phase 2 now collects.",
    })

    # ----- Slide 6  (was 13: Tech topic — Testing) ---------------------------
    # Existing copy is accurate per test_nightly.py + nightly.yml. Minor tweak.
    s = prs.slides[6]
    replace_in_runs(s, {
        "Isolation: google.cloud + mlflow stubbed via MagicMock — zero live cost in CI":
            "Isolation: google.cloud + mlflow stubbed via MagicMock — zero GCP cost in CI",
    })

    # ----- Slide 7  (was 14: Differentiators) --------------------------------
    s = prs.slides[7]
    replace_in_runs(s, {
        'Why this isn\'t just "ask ChatGPT."':
            'Why this isn\'t just "ask ChatGPT."  (Phase 2 shipped)',
        "Per-book check-ins reoptimize the rest of the curriculum — impossible from a single chat query.":
            "Per-book check-ins (status / rating / difficulty / comment) write to SQLite and gate future recommendations.",
        "Cloud Logging stores prior outputs and feeds them back into future curricula.":
            "Finished + abandoned books auto-exclude from future curricula. The feedback loop is closed.",
        "Two narrowly-scoped Gemini Flash calls + caching keep us at fractions of a cent per request.":
            "Two narrowly-scoped Gemini Flash calls + (user, survey) cache. ~$0.0003 per request.",
    })

    # ----- Slide 8  (was 15: Demo Video) -------------------------------------
    # Keep as-is; URL is in the closing slide.

    # ----- Slide 9  (was 16: Repo / Thank You) -------------------------------
    s = prs.slides[9]
    replace_in_runs(s, {
        "pipeline.py  ·  candidates.py  ·  curriculum.py":
            "pipeline.py · candidates.py · curriculum.py · reddit_extract.py · db.py",
    })

    prs.save(OUT)
    print(f"Wrote {OUT}")
    print(f"Slide count: {len(prs.slides)}")


if __name__ == "__main__":
    main()
