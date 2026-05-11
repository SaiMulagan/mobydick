"""
Reflect Sai's real MLflow + Docker + Cloud Run wiring (commit 8c161e1) on
slides 3 (Stack) and 6 (Models) of MobyDicks_Tech_v2.pptx.

What changed in the repo:
- mlflow-server/Dockerfile: tracking server on Cloud Run, Postgres + GCS backend
- train.py: TF-IDF recommender, logs to MLflow, registers BookRecommender, promotes to Production
- app.py: FastAPI service loads from registry, falls back to local copy
- Dockerfile: two-image setup (app + mlflow-server); model baked at build
"""

from pptx import Presentation

DECK = "/Users/aatishlobo/Desktop/mobydick/MobyDicks_Tech_v2.pptx"


def patch(slide, mapping):
    """Replace whole-run text matches; preserves font/size/color."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text in mapping:
                    run.text = mapping[run.text]


def main():
    prs = Presentation(DECK)

    # Slide 3 — Stack (0-indexed: slide index 2)
    s = prs.slides[2]
    patch(s, {
        # Docker: now two real images
        "Container image targeting Cloud Run; build script in repo root.":
            "Two images: app (FastAPI + baked model) and mlflow-server. Deployed to Cloud Run.",
        # MLflow: real tracking server, real registry
        "Run tracking stubbed in CI (MagicMock); planned for Phase 3 reranker.":
            "Cloud Run tracking server (Cloud SQL + GCS artifacts). Registers BookRecommender, promotes to Production.",
        # Cloud Run row: now two real services
        "Scale-to-zero target; ADC for Vertex AI + BigQuery + Cloud Logging.":
            "Two services live: app + mlflow-server. DB creds via Secret Manager, IAM for inter-service auth.",
    })

    # Slide 6 — Models (0-indexed: slide index 5)
    s = prs.slides[5]
    patch(s, {
        # Final bullet — make the "Phase 3 reranker" line concrete now that
        # the BookRecommender model is shipped.
        "Bonus: a third offline Gemini call extracts book mentions from Reddit posts. Phase 3 will train an MLflow-tracked reranker on the check-in data Phase 2 now collects.":
            "Bonus: a third offline Gemini call extracts book mentions from Reddit. Separate TF-IDF BookRecommender (train.py) logs to MLflow, registers + promotes to Production, served via FastAPI on Cloud Run.",
    })

    prs.save(DECK)
    print(f"Updated {DECK}")


if __name__ == "__main__":
    main()
